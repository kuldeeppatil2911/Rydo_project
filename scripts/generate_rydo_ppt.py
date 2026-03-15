from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


SLIDE_W = 13.333
SLIDE_H = 7.5

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_BG = RGBColor(9, 31, 74)
MID_BLUE = RGBColor(26, 91, 164)
CYAN = RGBColor(59, 180, 219)
LIME = RGBColor(176, 216, 34)
NAVY = RGBColor(19, 45, 79)
TEAL = RGBColor(22, 146, 153)
ORANGE = RGBColor(242, 139, 65)
PURPLE = RGBColor(111, 95, 208)
GREEN = RGBColor(45, 161, 110)
BG = RGBColor(244, 247, 250)
TEXT = RGBColor(43, 54, 66)
MUTED = RGBColor(94, 110, 126)
LINE = RGBColor(194, 204, 214)
PALE = RGBColor(232, 241, 248)
CARD = RGBColor(251, 253, 255)

INSTITUTE = "Parul Institute of Technology"
DEPARTMENT = "CSE Core"


def set_plain_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_full_rect(slide, x, y, w, h, color, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color
    return shape


def add_round_rect(slide, x, y, w, h, color, line=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Aptos",
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return box


def add_slide_header(slide, title, subtitle):
    add_full_rect(slide, 0, 0, SLIDE_W, 0.12, BLACK)
    add_full_rect(slide, 0, 0.12, SLIDE_W, 1.02, MID_BLUE)
    add_full_rect(slide, 12.0, 0.12, 0.78, 0.98, LIME)
    ribbon = add_round_rect(slide, 0.55, 0.28, 5.9, 0.55, CYAN, CYAN)
    ribbon.fill.transparency = 0.72
    ribbon.line.transparency = 1
    add_full_rect(slide, 0, 1.14, SLIDE_W, 0.06, CYAN)
    add_text(slide, 0.62, 0.3, 10.8, 0.38, title, 25, WHITE, True)
    add_text(slide, 0.66, 0.78, 11.0, 0.22, subtitle, 10.8, WHITE)


def add_colored_box(slide, x, y, w, h, text, color, size=15):
    shape = add_round_rect(slide, x, y, w, h, color)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Aptos"
    return shape


def add_outline_card(slide, x, y, w, h, title, lines, accent=TEAL):
    shape = add_round_rect(slide, x, y, w, h, CARD, accent)
    shape.line.width = Pt(1.6)
    shape.shadow.inherit = False
    add_text(slide, x + 0.12, y + 0.07, w - 0.24, 0.28, title, 14, accent, True)
    content = "\n".join(lines)
    add_text(slide, x + 0.14, y + 0.36, w - 0.28, h - 0.44, content, 10.5, TEXT)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=LINE, width=2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    connector.line.end_arrowhead = True
    return connector


def add_actor(slide, x, y, scale=1.0, color=NAVY):
    head = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.42 * scale), Inches(y), Inches(0.35 * scale), Inches(0.35 * scale)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = WHITE
    head.line.color.rgb = color
    head.line.width = Pt(1.8)

    def stick(x1, y1, x2, y2):
        c = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        c.line.color.rgb = color
        c.line.width = Pt(2)

    stick(x + 0.59 * scale, y + 0.35 * scale, x + 0.59 * scale, y + 0.95 * scale)
    stick(x + 0.22 * scale, y + 0.58 * scale, x + 0.96 * scale, y + 0.58 * scale)
    stick(x + 0.59 * scale, y + 0.95 * scale, x + 0.26 * scale, y + 1.33 * scale)
    stick(x + 0.59 * scale, y + 0.95 * scale, x + 0.92 * scale, y + 1.33 * scale)


def add_use_case(slide, x, y, w, h, text, color=TEAL):
    oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    oval.fill.solid()
    oval.fill.fore_color.rgb = WHITE
    oval.line.color.rgb = color
    oval.line.width = Pt(1.6)
    tf = oval.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11.5)
    r.font.bold = True
    r.font.color.rgb = TEXT
    r.font.name = "Aptos"
    return oval


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide, DARK_BG)

    add_full_rect(slide, 0, 0, SLIDE_W, 0.25, BLACK)
    add_full_rect(slide, 0, 0.25, SLIDE_W, SLIDE_H - 0.25, MID_BLUE)
    add_full_rect(slide, 0.0, 0.25, 13.333, 7.25, MID_BLUE)
    add_full_rect(slide, 12.0, 0.25, 0.78, 1.25, LIME)

    center_panel = add_round_rect(slide, 1.2, 1.1, 10.95, 4.8, CYAN, CYAN)
    center_panel.fill.transparency = 0.18
    center_panel.line.transparency = 1

    add_text(slide, 2.35, 1.95, 8.8, 0.9, "RYDO", 48, WHITE, True, PP_ALIGN.CENTER)
    add_full_rect(slide, 3.15, 3.08, 6.7, 0.5, BLACK)
    add_text(slide, 3.2, 3.12, 6.6, 0.3, "SYSTEM DESIGN & FLOW", 18, LIME, True, PP_ALIGN.CENTER)

    add_text(slide, 0.95, 5.15, 4.4, 0.35, "Minor Project - Project Definition", 15, WHITE, True)
    add_text(slide, 0.95, 5.72, 2.1, 0.3, "Presented By:", 15, WHITE, True)
    add_text(slide, 0.95, 6.03, 4.8, 0.7, "Kuldeep Ravindra Patil", 16, WHITE)

    add_text(slide, 9.8, 5.15, 2.0, 0.3, "Mentor:", 15, WHITE, True)
    add_text(slide, 9.8, 5.45, 3.0, 0.45, "Mr. Tadikonda Venkatata Durga Prasad", 14, WHITE)
    add_text(slide, 9.8, 5.95, 2.0, 0.3, "Institute:", 15, WHITE, True)
    add_text(slide, 9.8, 6.2, 2.9, 0.32, INSTITUTE, 13, WHITE)
    add_text(slide, 9.8, 6.48, 2.0, 0.3, "Department:", 15, WHITE, True)
    add_text(slide, 9.8, 6.73, 2.9, 0.28, DEPARTMENT, 13, WHITE)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "PROJECT OVERVIEW", "Problem statement, objective, and key innovation of the Rydo platform")

    add_outline_card(slide, 0.6, 1.35, 3.85, 2.3, "Problem Statement", [
        "Conventional demo ride-booking systems mainly focus on booking only.",
        "They usually do not include trusted-contact safety communication.",
        "Users need a simple platform for booking, tracking, and safety alerts in one flow.",
    ], NAVY)
    add_outline_card(slide, 4.75, 1.35, 3.85, 2.3, "Project Objective", [
        "Build a complete MERN stack web application for ride booking.",
        "Provide login, booking, live tracking, and emergency-contact support.",
        "Demonstrate real full-stack deployment with cloud database connectivity.",
    ], TEAL)
    add_outline_card(slide, 8.9, 1.35, 3.85, 2.3, "Key Innovation", [
        "Emergency alert can share vehicle number, driver, fare, payment mode, and OTP.",
        "Current-location map and moving-driver simulation improve demo quality.",
        "The system is presentation-ready and can be extended to a real product later.",
    ], ORANGE)

    add_colored_box(slide, 1.0, 4.3, 2.0, 0.85, "Authentication", NAVY)
    add_colored_box(slide, 3.35, 4.3, 2.0, 0.85, "Ride Booking", TEAL)
    add_colored_box(slide, 5.7, 4.3, 2.0, 0.85, "Tracking", PURPLE)
    add_colored_box(slide, 8.05, 4.3, 2.0, 0.85, "Safety Alert", ORANGE)
    add_colored_box(slide, 10.4, 4.3, 2.0, 0.85, "Cloud Deploy", GREEN)

    add_arrow(slide, 3.0, 4.73, 3.35, 4.73, TEAL)
    add_arrow(slide, 5.35, 4.73, 5.7, 4.73, TEAL)
    add_arrow(slide, 7.7, 4.73, 8.05, 4.73, TEAL)
    add_arrow(slide, 10.05, 4.73, 10.4, 4.73, TEAL)

    add_text(slide, 1.0, 5.45, 11.5, 0.75,
             "Core Result: Rydo combines secure login, ride estimation, live tracking, and emergency communication in a single user-friendly workflow.",
             15, NAVY, True, PP_ALIGN.CENTER)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "SYSTEM ARCHITECTURE", "High-level architecture showing frontend, backend, database, and external services")

    add_colored_box(slide, 0.65, 1.55, 1.65, 0.9, "User", NAVY, 17)
    add_colored_box(slide, 2.7, 1.2, 2.6, 1.3, "React Frontend\nLogin | Signup | Dashboard\nBook | Track | Profile", TEAL, 15)
    add_colored_box(slide, 5.85, 1.2, 2.2, 1.3, "Express Backend\nREST API + JWT", ORANGE, 16)
    add_colored_box(slide, 8.55, 0.85, 2.1, 0.95, "MongoDB Atlas", GREEN, 17)
    add_colored_box(slide, 8.55, 2.05, 2.1, 0.95, "SMTP / Twilio", PURPLE, 16)
    add_colored_box(slide, 11.05, 1.2, 1.7, 1.3, "Vercel +\nRailway", NAVY, 16)

    add_arrow(slide, 2.3, 1.98, 2.7, 1.85)
    add_arrow(slide, 5.3, 1.85, 5.85, 1.85)
    add_arrow(slide, 8.05, 1.55, 8.55, 1.32)
    add_arrow(slide, 8.05, 2.05, 8.55, 2.52)
    add_arrow(slide, 10.65, 1.85, 11.05, 1.85)

    add_outline_card(slide, 0.8, 3.45, 2.35, 2.25, "Frontend Layer", [
        "Built with React and Vite",
        "Uses React Router for page navigation",
        "AuthContext and RideContext manage state",
        "Responsive pages for login, booking, tracking, and profile",
    ], TEAL)
    add_outline_card(slide, 3.35, 3.45, 2.35, 2.25, "API Layer", [
        "Express routes expose business functions",
        "Controllers handle auth, booking, driver, meta, and user updates",
        "JWT protects private APIs",
        "CORS connects deployed frontend and backend",
    ], ORANGE)
    add_outline_card(slide, 5.9, 3.45, 2.35, 2.25, "Database Layer", [
        "MongoDB stores users, drivers, and bookings",
        "Mongoose models enforce schema structure",
        "Seed data helps simulate drivers and ride options",
        "Recent bookings are fetched for dashboard visibility",
    ], GREEN)
    add_outline_card(slide, 8.45, 3.45, 2.35, 2.25, "Safety Layer", [
        "Emergency contact stored in user profile",
        "Alert includes vehicle number, driver, fare, and OTP",
        "Email works through SMTP",
        "SMS works through Twilio integration",
    ], PURPLE)
    add_outline_card(slide, 11.0, 3.45, 1.9, 2.25, "Deployment", [
        "Frontend deployed on Vercel",
        "Backend deployed on Railway",
        "Database hosted on Atlas",
        "Public links available for demo",
    ], NAVY)


def slide_modules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "MODULE DESIGN AND PROCESS FLOW", "Detailed breakdown of each major module inside the Rydo application")

    add_outline_card(slide, 0.65, 1.35, 2.35, 2.3, "1. Authentication", [
        "Register new user",
        "Login existing user",
        "Password hashing using bcryptjs",
        "JWT token generation and validation",
    ], NAVY)
    add_outline_card(slide, 3.2, 1.35, 2.35, 2.3, "2. Profile", [
        "Store phone number",
        "Store emergency contact details",
        "Enable or disable emergency alerts",
        "Load current user profile after login",
    ], TEAL)
    add_outline_card(slide, 5.75, 1.35, 2.35, 2.3, "3. Ride Booking", [
        "Select pickup and drop points",
        "Choose ride type and payment mode",
        "Estimate distance, time, and fare",
        "Create booking with OTP and status",
    ], ORANGE)
    add_outline_card(slide, 8.3, 1.35, 2.35, 2.3, "4. Driver Allocation", [
        "Filter available drivers",
        "Match driver to selected ride type",
        "Attach driver snapshot to booking",
        "Prepare ride details for tracking",
    ], PURPLE)
    add_outline_card(slide, 10.85, 1.35, 1.85, 2.3, "5. Alerts", [
        "Check emergency contact",
        "Prepare email and SMS body",
        "Include vehicle, fare, and OTP",
        "Send via SMTP/Twilio when configured",
    ], GREEN)

    add_colored_box(slide, 0.95, 4.4, 2.0, 0.85, "Signup / Login", NAVY)
    add_colored_box(slide, 3.2, 4.4, 2.0, 0.85, "Profile Ready", TEAL)
    add_colored_box(slide, 5.45, 4.4, 2.0, 0.85, "Book Ride", ORANGE)
    add_colored_box(slide, 7.7, 4.4, 2.0, 0.85, "Track Ride", PURPLE)
    add_colored_box(slide, 9.95, 4.4, 2.0, 0.85, "Send Alert", GREEN)

    add_arrow(slide, 2.95, 4.82, 3.2, 4.82, NAVY)
    add_arrow(slide, 5.2, 4.82, 5.45, 4.82, NAVY)
    add_arrow(slide, 7.45, 4.82, 7.7, 4.82, NAVY)
    add_arrow(slide, 9.7, 4.82, 9.95, 4.82, NAVY)

    add_text(slide, 1.1, 5.55, 11.0, 0.65,
             "This sequential flow makes the UI simple for users: first authenticate, then book, then track, and finally notify the trusted contact if safety mode is enabled.",
             14, TEXT, False, PP_ALIGN.CENTER)


def slide_dfd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "DFD (DATA FLOW DIAGRAM)", "Logical movement of ride, user, and safety data across the system")

    add_colored_box(slide, 0.55, 2.3, 1.55, 0.9, "User", NAVY, 18)
    add_colored_box(slide, 2.6, 1.1, 2.25, 0.8, "Authentication Process", TEAL, 14)
    add_colored_box(slide, 2.6, 2.3, 2.25, 0.8, "Ride Booking Process", ORANGE, 14)
    add_colored_box(slide, 2.6, 3.5, 2.25, 0.8, "Tracking Process", PURPLE, 14)
    add_colored_box(slide, 5.45, 0.8, 2.15, 0.8, "User Data Store", GREEN, 14)
    add_colored_box(slide, 5.45, 2.0, 2.15, 0.8, "Booking Data Store", GREEN, 14)
    add_colored_box(slide, 5.45, 3.2, 2.15, 0.8, "Driver Data Store", GREEN, 14)
    add_colored_box(slide, 8.2, 1.2, 2.55, 0.9, "Emergency Alert Service", NAVY, 15)
    add_colored_box(slide, 8.2, 3.05, 2.55, 0.9, "Map and Location Service", TEAL, 15)
    add_colored_box(slide, 11.0, 2.3, 1.8, 0.9, "Emergency Contact", ORANGE, 14)

    add_arrow(slide, 2.1, 2.75, 2.6, 1.5)
    add_arrow(slide, 2.1, 2.75, 2.6, 2.7)
    add_arrow(slide, 2.1, 2.75, 2.6, 3.9)

    add_arrow(slide, 4.85, 1.5, 5.45, 1.2)
    add_arrow(slide, 4.85, 2.7, 5.45, 2.4)
    add_arrow(slide, 4.85, 3.9, 5.45, 3.6)

    add_arrow(slide, 7.6, 1.2, 8.2, 1.65)
    add_arrow(slide, 7.6, 2.4, 8.2, 1.75)
    add_arrow(slide, 7.6, 3.6, 8.2, 3.5)
    add_arrow(slide, 10.75, 1.7, 11.0, 2.6)
    add_arrow(slide, 10.75, 3.5, 11.0, 2.75)

    add_text(slide, 1.86, 1.48, 0.7, 0.2, "login", 9, MUTED)
    add_text(slide, 1.72, 2.45, 0.95, 0.2, "booking request", 9, MUTED)
    add_text(slide, 1.82, 3.68, 0.85, 0.2, "track ride", 9, MUTED)
    add_text(slide, 8.62, 2.2, 1.75, 0.25, "alert details sent", 9, MUTED, False, PP_ALIGN.CENTER)


def slide_er(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "ER DIAGRAM", "Entity relationships implemented in MongoDB through Mongoose models")

    add_outline_card(slide, 0.55, 1.35, 2.65, 2.95, "USER", [
        "userId (PK)",
        "name",
        "email",
        "password",
        "phone",
        "emergencyAlertsEnabled",
        "emergencyContact.name",
        "emergencyContact.phone",
        "emergencyContact.email",
    ], NAVY)
    add_outline_card(slide, 3.55, 1.15, 2.95, 3.35, "BOOKING", [
        "bookingId (PK)",
        "userId (FK)",
        "pickup",
        "dropoff",
        "rideType",
        "driver snapshot",
        "payment",
        "tripMode",
        "distance, time, fare",
        "otp, status, activityLog",
    ], TEAL)
    add_outline_card(slide, 6.85, 1.35, 2.45, 2.95, "DRIVER", [
        "driverId (PK)",
        "name",
        "vehicle",
        "plate",
        "rating",
        "eta",
        "rideType",
        "status",
    ], ORANGE)
    add_outline_card(slide, 9.65, 1.35, 2.95, 2.95, "LOCATION / META", [
        "locationId (PK)",
        "name",
        "x coordinate",
        "y coordinate",
        "ride type options",
        "base estimate values",
    ], PURPLE)

    add_arrow(slide, 3.2, 2.75, 3.55, 2.75)
    add_arrow(slide, 6.5, 2.75, 6.85, 2.75)
    add_arrow(slide, 9.3, 2.75, 9.65, 2.75)

    add_text(slide, 3.18, 2.42, 0.7, 0.2, "1 : M", 10, NAVY, True)
    add_text(slide, 6.48, 2.42, 0.7, 0.2, "M : 1", 10, NAVY, True)
    add_text(slide, 9.28, 2.42, 0.7, 0.2, "M : 1", 10, NAVY, True)

    add_outline_card(slide, 1.25, 5.0, 10.9, 0.95, "Relationship Summary", [
        "A single user can create many bookings. Every booking stores ride metadata and a selected driver snapshot. Locations and ride-type data are referenced during estimate and booking generation.",
    ], GREEN)


def slide_use_case(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "USE CASE DIAGRAM", "Main interactions between system actors and Rydo features")

    system = add_round_rect(slide, 2.15, 1.0, 9.0, 5.7, WHITE, LINE)
    system.line.width = Pt(1.5)
    add_text(slide, 5.3, 1.15, 2.7, 0.25, "RYDO SYSTEM", 16, NAVY, True, PP_ALIGN.CENTER)

    add_actor(slide, 0.55, 2.2, 1.0, NAVY)
    add_text(slide, 0.28, 3.68, 1.35, 0.25, "User", 13, NAVY, True, PP_ALIGN.CENTER)
    add_actor(slide, 11.55, 2.2, 1.0, ORANGE)
    add_text(slide, 11.0, 3.68, 2.4, 0.25, "Emergency Contact", 13, ORANGE, True, PP_ALIGN.CENTER)

    points = [
        (3.0, 1.8, 1.9, 0.76, "Register / Login", NAVY),
        (5.35, 1.8, 1.9, 0.76, "Manage Profile", TEAL),
        (7.7, 1.8, 2.0, 0.76, "Set Emergency Contact", ORANGE),
        (3.0, 3.0, 1.9, 0.76, "Book Ride", TEAL),
        (5.35, 3.0, 1.9, 0.76, "View Fare Estimate", PURPLE),
        (7.7, 3.0, 2.0, 0.76, "Track Ride", GREEN),
        (4.1, 4.35, 2.2, 0.76, "Receive OTP", NAVY),
        (6.95, 4.35, 2.1, 0.76, "Receive Alert", ORANGE),
    ]

    for x, y, w, h, text, color in points:
        add_use_case(slide, x, y, w, h, text, color)

    for end_x, end_y in [(3.0, 2.18), (5.35, 2.18), (7.7, 2.18), (3.0, 3.38), (5.35, 3.38), (7.7, 3.38), (4.1, 4.73)]:
        c = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, Inches(1.58), Inches(2.82), Inches(end_x), Inches(end_y)
        )
        c.line.color.rgb = LINE
        c.line.width = Pt(1.4)

    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, Inches(11.55), Inches(2.82), Inches(9.05), Inches(4.73)
    )
    c.line.color.rgb = LINE
    c.line.width = Pt(1.4)


def slide_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "TECH STACK AND DEPLOYMENT", "Tools, libraries, and hosting platforms used to implement the project")

    add_outline_card(slide, 0.65, 1.35, 2.4, 2.35, "Frontend Stack", [
        "React.js",
        "Vite",
        "React Router",
        "React Leaflet + Leaflet",
        "Context API for state",
    ], TEAL)
    add_outline_card(slide, 3.3, 1.35, 2.4, 2.35, "Backend Stack", [
        "Node.js",
        "Express.js",
        "JWT authentication",
        "bcryptjs for hashing",
        "REST API architecture",
    ], ORANGE)
    add_outline_card(slide, 5.95, 1.35, 2.4, 2.35, "Database Stack", [
        "MongoDB Atlas",
        "Mongoose ODM",
        "User collection",
        "Booking collection",
        "Driver collection",
    ], GREEN)
    add_outline_card(slide, 8.6, 1.35, 2.4, 2.35, "Notification Stack", [
        "Nodemailer",
        "SMTP integration",
        "Twilio SMS support",
        "Emergency safety workflow",
        "Cloud environment variables",
    ], PURPLE)
    add_outline_card(slide, 11.25, 1.35, 1.45, 2.35, "Hosting", [
        "Vercel",
        "Railway",
        "Atlas",
        "GitHub",
    ], NAVY)

    add_colored_box(slide, 1.2, 4.35, 2.0, 0.85, "GitHub Repo", NAVY)
    add_colored_box(slide, 4.0, 4.35, 2.0, 0.85, "Railway API", ORANGE)
    add_colored_box(slide, 6.8, 4.35, 2.0, 0.85, "Vercel Client", TEAL)
    add_colored_box(slide, 9.6, 4.35, 2.0, 0.85, "Atlas Database", GREEN)

    add_arrow(slide, 3.2, 4.78, 4.0, 4.78, NAVY)
    add_arrow(slide, 6.0, 4.78, 6.8, 4.78, NAVY)
    add_arrow(slide, 8.8, 4.78, 9.6, 4.78, NAVY)

    add_text(slide, 0.95, 5.55, 11.4, 0.7,
             "Deployment Outcome: the application is hosted online, allowing real login, booking, and API communication for project demonstration.",
             14, NAVY, True, PP_ALIGN.CENTER)


def slide_screenshots(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "APPLICATION SCREENSHOTS", "Placeholder frames added so you can quickly insert live project screenshots before final submission")

    add_outline_card(slide, 0.65, 1.25, 3.95, 4.65, "Login / Signup Screen", [
        "Replace this area with a screenshot of the login or signup page.",
        "Recommended capture: app logo, email/password form, and clean UI.",
    ], NAVY)
    add_outline_card(slide, 4.7, 1.25, 3.95, 4.65, "Ride Booking Screen", [
        "Replace this area with the booking page screenshot.",
        "Recommended capture: pickup, drop, ride type, fare estimate, and CTA button.",
    ], TEAL)
    add_outline_card(slide, 8.75, 1.25, 3.95, 4.65, "Tracking / Profile Screen", [
        "Replace this area with a screenshot of live tracking or emergency-contact profile.",
        "Recommended capture: map, driver details, or safety toggle section.",
    ], ORANGE)

    add_text(slide, 1.0, 3.15, 3.25, 0.7, "Insert Screenshot Here", 18, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, 5.05, 3.15, 3.25, 0.7, "Insert Screenshot Here", 18, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, 9.1, 3.15, 3.25, 0.7, "Insert Screenshot Here", 18, MUTED, True, PP_ALIGN.CENTER)

    add_text(slide, 1.1, 6.1, 11.1, 0.45,
             "Tip: Open the live app, take three screenshots, and replace these placeholders directly in PowerPoint.",
             12, MUTED, False, PP_ALIGN.CENTER)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide)
    add_slide_header(slide, "CONCLUSION AND FUTURE SCOPE", "Summary of outcomes and possible next improvements")

    add_outline_card(slide, 0.7, 1.4, 5.75, 3.85, "Project Conclusion", [
        "Rydo successfully demonstrates a complete MERN stack ride-booking workflow.",
        "The project covers user authentication, booking, tracking, profile management, and emergency safety alerts.",
        "It combines academic concepts such as database design, REST APIs, cloud deployment, and responsive UI in one solution.",
        "The deployed version proves that the system can run as a real web application, not only as a local prototype.",
    ], NAVY)
    add_outline_card(slide, 6.75, 1.4, 5.85, 3.85, "Future Scope", [
        "Integrate real-time GPS route navigation and accurate map routing.",
        "Add online payment gateway and wallet support.",
        "Create an admin dashboard for driver and booking management.",
        "Add ride scheduling, push notifications, and multilingual support.",
        "Improve analytics, traffic-aware ETA prediction, and real production messaging.",
    ], TEAL)

    add_text(slide, 2.35, 6.15, 8.5, 0.35, "Presented By: Kuldeep Ravindra Patil | Mentor: Mr. Tadikonda Venkatata Durga Prasad", 14, MUTED, False, PP_ALIGN.CENTER)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_plain_bg(slide, MID_BLUE)

    hero = add_round_rect(slide, 1.1, 1.15, 11.1, 5.1, CYAN, CYAN)
    hero.fill.transparency = 0.18
    hero.line.transparency = 1
    add_full_rect(slide, 0.0, 0.0, 13.333, 0.22, BLACK)
    add_full_rect(slide, 12.0, 0.22, 0.78, 1.18, LIME)
    add_text(slide, 2.45, 2.15, 8.3, 0.8, "THANK YOU", 38, WHITE, True, PP_ALIGN.CENTER)
    add_full_rect(slide, 3.3, 3.2, 6.7, 0.52, BLACK)
    add_text(slide, 3.35, 3.25, 6.6, 0.3, "Questions and Discussion", 18, LIME, True, PP_ALIGN.CENTER)
    add_text(slide, 2.2, 4.35, 9.0, 0.35, "Presented By: Kuldeep Ravindra Patil", 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 2.2, 4.8, 9.0, 0.35, "Mentor: Mr. Tadikonda Venkatata Durga Prasad", 16, WHITE, False, PP_ALIGN.CENTER)
    add_text(slide, 2.2, 5.25, 9.0, 0.35, f"Institute: {INSTITUTE} | Department: {DEPARTMENT}", 14, WHITE, False, PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_cover(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_modules(prs)
    slide_dfd(prs)
    slide_er(prs)
    slide_use_case(prs)
    slide_tech_stack(prs)
    slide_screenshots(prs)
    slide_conclusion(prs)
    slide_thank_you(prs)

    output = "Rydo_Detailed_System_Design_Flow.pptx"
    prs.save(output)
    print(output)


if __name__ == "__main__":
    main()
